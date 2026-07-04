import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0] # Title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "E-Commerce Product Catalog System"
subtitle.text = "A Java Console Application demonstrating advanced OOP Principles\n\nPresenter: [Your Name]\nDate: [Date]"

# Helper to add bullet slides
def add_bullet_slide(title_text, bullets):
    layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    
    tf = slide.placeholders[1].text_frame
    tf.text = bullets[0]
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0

# Slide 2
add_bullet_slide(
    "Project Overview & The Business Problem",
    [
        "The Problem: Modern platforms sell physical goods, digital downloads, and subscriptions, each with different rules.",
        "The Goal: Build an extensible catalog system that avoids messy if/else or instanceof checks.",
        "The Solution: A robust Java application utilizing a polymorphic product hierarchy to dynamically handle diverse business logic."
    ]
)

# Slide 3
add_bullet_slide(
    "Core Technologies & Architecture",
    [
        "Language: Java",
        "Paradigm: Object-Oriented Programming (OOP)",
        "Data Persistence: CSV File I/O",
        "Design Patterns Used: Template Method Pattern & Factory-style Parsing",
        "UI: Interactive Command Line Interface (CLI)"
    ]
)

# Slide 4
add_bullet_slide(
    "Object-Oriented Design (OOP)",
    [
        "Abstract Base Class (Product.java): Holds shared properties (ID, Name, Base Price).",
        "Abstract Hooks: Defines methods children must implement (calculateShippingCost, calculateTax).",
        "Subclasses: PhysicalProduct, DigitalProduct, SubscriptionProduct inherit and implement unique logic.",
        "Benefit: CatalogService processes thousands of mixed products polymorphically."
    ]
)

# Slide 5
add_bullet_slide(
    "Data Persistence & File I/O",
    [
        "The Challenge: Storing polymorphic objects into a flat text file.",
        "Implementation: Dedicated Repository classes (ProductRepository, OrderRepository).",
        "Each product is saved with a Type column (PHYSICAL, DIGITAL).",
        "On startup, the system reads the CSV and seamlessly resurrects the correct child class."
    ]
)

# Slide 6
add_bullet_slide(
    "Dynamic Business Logic",
    [
        "Real-Time Cost Calculation: Costs react dynamically to a 'place' parameter (e.g., Texas, Hawaii).",
        "Simulated Logistics: International string-length algorithms and remote keyword detection.",
        "Auto-ID Generation: System automatically parses CSV states to generate sequential IDs (e.g., P-001) mimicking SQL."
    ]
)

# Slide 7
add_bullet_slide(
    "Security & User Flows",
    [
        "Role-Based Access Control (RBAC): Users route based on role (ADMIN or CUSTOMER).",
        "Admin Portal: View specific catalogs and add new products.",
        "Customer Portal: View simplified catalog, simulate purchases.",
        "Data Privacy: Securely track order history restricted to the logged-in user."
    ]
)

# Slide 8
add_bullet_slide(
    "Conclusion & Future Enhancements",
    [
        "Summary: Successfully built a scalable, OOP-driven console app mirroring real-world business logic.",
        "Future 1: Migrate from CSV storage to a relational database (MySQL).",
        "Future 2: Implement a GUI using JavaFX or a web frontend.",
        "Future 3: Add a multi-product shopping cart feature.",
        "Q&A"
    ]
)

# Save
prs.save("Ecommerce_Presentation.pptx")
print("Presentation generated successfully: Ecommerce_Presentation.pptx")
